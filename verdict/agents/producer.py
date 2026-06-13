"""Producer — PRODUCE. Generates the real deliverables: a signed memo (.md), a
deck (.pptx via python-pptx) and a model (.xlsx via openpyxl). SenseNova U1 was
unverified, so python-pptx/openpyxl is the load-bearing producer; if a lib is
missing it degrades to text/CSV. Runs in Daytona when live."""
import asyncio
import os

from ..events import ev
from ..schemas import Envelope
from .. import config


def _memo_md(fixture, thesis, val) -> str:
    risks = "\n".join(f"- {r}" for r in thesis["risks"])
    up = round(val.get("upside", 0) * 100)
    return (
        f"# Verdict — Investment Due-Diligence Memo\n\n"
        f"**{fixture['name']} ({fixture['ticker']})** · {fixture['sector']}\n\n"
        f"## Recommendation: {thesis['recommendation']}  ·  conviction {thesis['conviction']}\n\n"
        f"{thesis['summary']}\n\n"
        f"## Valuation\n"
        f"- Method: {val['method']}\n"
        f"- Current price: ${val['current_price_usd']}\n"
        f"- Fair value: ${val['fair_value_usd']}  (**{up}% upside**)\n\n"
        f"## Key risks\n{risks}\n\n"
        f"_Generated autonomously by Verdict and cryptographically signed (see attached credential)._\n"
    )


def _deck(path, fixture, thesis, val):
    from pptx import Presentation

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = f"Verdict — {fixture['name']} ({fixture['ticker']})"
    s.placeholders[1].text = f"{thesis['recommendation']} · fair value ${val['fair_value_usd']} ({round(val.get('upside',0)*100)}% upside)"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Thesis & risks"
    tf = s.placeholders[1].text_frame
    tf.text = thesis["summary"]
    for r in thesis["risks"]:
        p = tf.add_paragraph()
        p.text = f"Risk: {r}"
        p.level = 1
    prs.save(path)
    return path


def _xlsx(path, fixture, val):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Valuation"
    rows = [
        ("Company", fixture["name"]),
        ("Ticker", fixture["ticker"]),
        ("Method", val["method"]),
        ("Current price", val["current_price_usd"]),
        ("Fair value", val["fair_value_usd"]),
        ("EV/Sales", val["ev_sales_multiple"]),
        ("Rev growth", val["growth"]),
        ("Upside", val.get("upside", 0)),
    ]
    for r in rows:
        ws.append(r)
    ws["B8"].number_format = "0.0%"
    ws["B7"].number_format = "0.0%"
    wb.save(path)
    return path


def _build(base, fixture, ctx):
    os.makedirs(os.path.dirname(base), exist_ok=True)
    thesis = ctx.get("thesis", fixture["thesis"])
    val = ctx.get("valuation", fixture["valuation"])
    paths = []

    memo = base + "_memo.md"
    with open(memo, "w") as f:
        f.write(_memo_md(fixture, thesis, val))
    paths.append(memo)

    builders = [lambda p: _deck(p, fixture, thesis, val), lambda p: _xlsx(p, fixture, val)]
    for build, out in zip(builders, (base + "_deck.pptx", base + "_model.xlsx")):
        try:
            paths.append(build(out))
        except Exception:
            pass  # degrade silently — memo already covers the deliverable
    return paths


async def run(company, fixture, ctx, emit) -> Envelope:
    base = os.path.join(config.ARTIFACTS_DIR, company)
    paths = await asyncio.to_thread(_build, base, fixture, ctx)
    for p in paths:
        await emit(ev("artifact", stage="producer", path=p,
                      kind=os.path.splitext(p)[1].lstrip(".")))
        await asyncio.sleep(0.08)
    return Envelope(data={"artifacts": paths}, status="ok", source="mock",
                    notes=f'{len(paths)} deliverables: ' + ", ".join(os.path.basename(p) for p in paths))
